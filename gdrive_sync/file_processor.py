# -*- coding: utf-8 -*-
"""
file_processor.py \u2014 Extracci\u00f3n de texto en memoria por tipo de archivo
"""

import io
import re
import json
import logging

logger = logging.getLogger(__name__)


def _expand_acronyms(text):
    """Detecta acrónimos con puntos (P.A.C.O.) y agrega versión sin puntos (PACO)."""
    def _add_plain(m):
        original = m.group(0)
        plain = original.replace(".", "")
        if plain == original:
            return original
        return f"{original} ({plain})"
    return re.sub(r'(?:[A-Z]\.){2,}[A-Z]?\.?', _add_plain, text)


def clean_text(text, source_type="generic"):
    """Limpia texto antes de chunking y vectorizaci\u00f3n."""
    if not text or not text.strip():
        return ""

    lines = text.splitlines()

    # Detectar headers/footers repetidos (comunes en PDFs corporativos)
    if len(lines) > 10:
        line_counts = {}
        for line in lines:
            clean_line = line.strip().lower()
            if clean_line and len(clean_line) < 100:
                line_counts[clean_line] = line_counts.get(clean_line, 0) + 1
        page_count = max(1, len(lines) // 40)
        repeated = {l for l, c in line_counts.items() if c > page_count * 0.3 and c > 2}
        lines = [l for l in lines if l.strip().lower() not in repeated]

    text = "\n".join(lines)

    # Quitar n\u00fameros de p\u00e1gina sueltos
    text = re.sub(r'\n\s*\d{1,3}\s*\n', '\n', text)
    text = re.sub(r'\n\s*(?:P\u00e1gina|Page)\s+\d+\s*(?:de|of)\s+\d+\s*\n', '\n', text, flags=re.IGNORECASE)

    # Normalizar whitespace excesivo
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)

    # Expandir acrónimos: P.A.C.O. -> P.A.C.O. (PACO)
    text = _expand_acronyms(text)

    return text.strip()


def extract_text(buffer: io.BytesIO, mime_type: str, filename: str) -> str:
    try:
        mime = mime_type.lower()
        if mime == "application/pdf":
            return _extract_pdf(buffer, filename)
        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_docx(buffer, filename)
        elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _extract_pptx(buffer, filename)
        elif mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return _extract_xlsx(buffer, filename)
        elif mime == "text/csv":
            return _extract_csv(buffer, filename)
        elif mime == "application/json":
            return _extract_json(buffer, filename)
        elif mime in ("text/plain", "text/markdown"):
            return _extract_text(buffer, filename)
        elif mime == "text/html":
            return _extract_html(buffer, filename)
        elif mime in ("image/jpeg", "image/png", "image/webp", "image/gif"):
            return _extract_image_vision(buffer, filename)
        else:
            logger.warning(f"   Formato no soportado: {mime} ({filename})")
            return ""
    except Exception as e:
        logger.error(f"   Error extrayendo texto de '{filename}': {e}")
        return ""


def extract_metadata(buffer: io.BytesIO, mime_type: str, filename: str) -> dict:
    """Extrae metadatos del archivo (t\u00edtulo, autor, fecha)."""
    meta = {"title": filename.rsplit(".", 1)[0], "author": "", "created_date": ""}
    try:
        mime = mime_type.lower()
        if mime == "application/pdf":
            import fitz
            pos = buffer.tell()
            doc = fitz.open(stream=buffer.read(), filetype="pdf")
            buffer.seek(pos)
            m = doc.metadata
            if m.get("title"):
                meta["title"] = m["title"]
            if m.get("author"):
                meta["author"] = m["author"]
            if m.get("creationDate"):
                meta["created_date"] = m["creationDate"]
            meta["page_count"] = len(doc)
            doc.close()
        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            pos = buffer.tell()
            doc = Document(buffer)
            buffer.seek(pos)
            props = doc.core_properties
            if props.title:
                meta["title"] = props.title
            if props.author:
                meta["author"] = props.author
            if props.created:
                meta["created_date"] = str(props.created)
    except Exception as e:
        logger.debug(f"   No se pudieron extraer metadatos de '{filename}': {e}")
    return meta


def _extract_pdf(buffer, filename):
    import fitz, io
    buffer.seek(0)
    raw_bytes = buffer.read()
    text_parts = []
    page_samples = []
    with fitz.open(stream=raw_bytes, filetype="pdf") as doc:
        for i, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                text_parts.append(text)
            if i < 3:
                page_samples.append(text.strip())
    raw = "\n\n".join(text_parts)
    result = clean_text(raw, "pdf")
    # Detectar PDF escaneado: pocas páginas con texto o texto total muy corto
    total_sample = " ".join(page_samples)
    if len(total_sample) < 100 or len(result) < 200:
        logger.info(f"   PDF escaneado detectado: '{filename}' — usando visión")
        return _extract_pdf_vision(raw_bytes, filename)
    logger.debug(f"   PDF '{filename}': {len(result)} caracteres")
    return result


def _extract_pdf_vision(raw_bytes: bytes, filename: str) -> str:
    """Convierte páginas de PDF escaneado a imágenes y usa visión — paralelo + resolución optimizada."""
    import fitz, base64, httpx, asyncio

    CONCURRENCY = 3       # páginas simultáneas
    MATRIX_SCALE = 1.5    # reducido de 2.0 → ~44% menos datos por imagen

    prompt = """Transcribe todo el texto visible en esta página manteniendo estructura y formato.
Si hay tablas, mantenlas con separadores |. Responde solo con el texto transcrito, sin comentarios."""

    async def _process_page(sem, i, total, img_b64):
        async with sem:
            async with httpx.AsyncClient(timeout=180) as client:
                resp = await client.post("http://localhost:11434/api/chat", json={
                    "model": "sonia-qwen:9b",
                    "think": False,
                    "stream": False,
                    "messages": [{"role": "user", "content": prompt, "images": [img_b64]}]
                })
                text = resp.json().get("message", {}).get("content", "").strip()
                logger.debug(f"   Página {i+1}/{total}: {len(text)} chars")
                return i, text

    async def _process_all(pages_b64):
        sem = asyncio.Semaphore(CONCURRENCY)
        tasks = [_process_page(sem, i, len(pages_b64), b64) for i, b64 in enumerate(pages_b64)]
        results = await asyncio.gather(*tasks)
        return sorted(results, key=lambda x: x[0])

    try:
        # Convertir todas las páginas a base64
        pages_b64 = []
        with fitz.open(stream=raw_bytes, filetype="pdf") as doc:
            total = len(doc)
            logger.info(f"   Procesando {total} páginas con visión paralela (x{CONCURRENCY}): {filename}")
            mat = fitz.Matrix(MATRIX_SCALE, MATRIX_SCALE)
            for page in doc:
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("jpeg")
                pages_b64.append(base64.b64encode(img_bytes).decode("utf-8"))

        # Procesar en paralelo
        results = asyncio.run(_process_all(pages_b64))

        texts = []
        for i, text in results:
            if text:
                texts.append(f"[Página {i+1}/{total}]\n{text}")

        result = "\n\n".join(texts)
        logger.info(f"   Vision PDF OK: {len(result)} chars — {filename}")
        return result
    except Exception as e:
        logger.error(f"   Error vision PDF {filename}: {e}")
        return ""

def _extract_docx(buffer, filename):
    from docx import Document
    doc = Document(buffer)
    paragraphs = []
    current_heading = ""
    for p in doc.paragraphs:
        if not p.text.strip():
            continue
        if p.style.name.startswith('Heading'):
            current_heading = p.text.strip()
            paragraphs.append(f"## {current_heading}")
        else:
            paragraphs.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    raw = "\n\n".join(paragraphs)
    result = clean_text(raw, "docx")
    logger.debug(f"   DOCX '{filename}': {len(result)} caracteres")
    return result

def _extract_pptx(buffer, filename):
    from pptx import Presentation
    prs = Presentation(buffer)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        # Extraer notas del presentador
        notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                notes = notes_frame.text.strip()
        slide_text = f"[Diapositiva {i+1}]\n" + "\n".join(parts)
        if notes:
            slide_text += f"\n[Notas]: {notes}"
        if parts or notes:
            slides_text.append(slide_text)
    result = "\n\n".join(slides_text)
    logger.debug(f"   PPTX '{filename}': {len(result)} caracteres")
    return result

def _extract_xlsx(buffer, filename):
    import openpyxl
    wb = openpyxl.load_workbook(buffer, read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            row_data = " | ".join(str(cell) for cell in row if cell is not None)
            if row_data.strip():
                rows_text.append(row_data)
        if rows_text:
            sheets_text.append(f"[Hoja: {sheet_name}]\n" + "\n".join(rows_text))
    wb.close()
    result = "\n\n".join(sheets_text)
    logger.debug(f"   XLSX '{filename}': {len(result)} caracteres")
    return result

def _extract_csv(buffer, filename):
    import csv
    content = buffer.read().decode("utf-8", errors="replace")
    reader = csv.reader(content.splitlines())
    rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
    result = "\n".join(rows)
    logger.debug(f"   CSV '{filename}': {len(result)} caracteres")
    return result

def _extract_json(buffer, filename):
    content = buffer.read().decode("utf-8", errors="replace")
    try:
        data = json.loads(content)
        result = json.dumps(data, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        result = content
    logger.debug(f"   JSON '{filename}': {len(result)} caracteres")
    return result

def _extract_html(buffer, filename):
    try:
        from bs4 import BeautifulSoup
        html = buffer.read().decode("utf-8", errors="replace")
        soup = BeautifulSoup(html, "lxml")
        for tag in soup(["script", "style", "head"]):
            tag.decompose()
        text = soup.get_text(separator=" ")
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        return " ".join(lines)
    except Exception as e:
        logger.error(f"Error HTML {filename}: {e}")
        return ""


def _extract_image_vision(buffer, filename):
    """Describe una imagen usando Qwen3.5 visi\u00f3n via Ollama."""
    import base64
    import httpx
    try:
        image_b64 = base64.b64encode(buffer.read()).decode("utf-8")
        prompt = """Analiza esta imagen con detalle:
1. Si contiene TEXTO: transcribe todo el texto visible manteniendo la estructura.
2. Si contiene DIAGRAMAS o GRAFICOS: describe los elementos clave, relaciones y datos.
3. Si es MIXTO: primero transcribe el texto, luego describe los elementos visuales.
Responde en espa\u00f1ol. Se exhaustivo."""

        with httpx.Client(timeout=120) as client:
            resp = client.post("http://localhost:11434/api/chat", json={
                "model":  "sonia-qwen:9b",
                "think":  False,
                "stream": False,
                "messages": [{"role": "user", "content": prompt, "images": [image_b64]}]
            })
            result = resp.json().get("message", {}).get("content", "").strip()
            logger.debug(f"   Vision OK: {len(result)} chars \u2014 {filename}")
            return result
    except Exception as e:
        logger.error(f"   Error vision {filename}: {e}")
        return ""


def _extract_text(buffer, filename):
    result = buffer.read().decode("utf-8", errors="replace")
    logger.debug(f"   TXT/MD '{filename}': {len(result)} caracteres")
    return result
