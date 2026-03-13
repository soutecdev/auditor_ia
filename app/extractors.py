# -*- coding: utf-8 -*-
"""Standalone text extraction functions for file upload.
No heavy deps (no SentenceTransformer, no Qdrant).
Copied from scripts/ingest.py to avoid module-level imports."""

import re
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def _expand_acronyms(text):
    """Detecta acr\u00f3nimos con puntos (P.A.C.O.) y agrega versi\u00f3n sin puntos."""
    def _add_plain(m):
        original = m.group(0)
        plain = original.replace(".", "")
        if plain == original:
            return original
        return f"{original} ({plain})"
    return re.sub(r'(?:[A-Z]\.){2,}[A-Z]?\.?', _add_plain, text)


def clean_text(text, source_type="generic"):
    if not text or not text.strip():
        return ""
    lines = text.splitlines()
    if len(lines) > 10:
        line_counts = {}
        for line in lines:
            clean_line = line.strip().lower()
            if clean_line and len(clean_line) < 100:
                line_counts[clean_line] = line_counts.get(clean_line, 0) + 1
        page_count = max(1, len(lines) // 40)
        repeated = {l for l, c in line_counts.items() if c > page_count * 0.3 and c > 2}
        lines = [l for l in lines if l.strip().lower() not in repeated]
    text = "\n".join(lines)
    text = re.sub(r'\n\s*\d{1,3}\s*\n', '\n', text)
    text = re.sub(r'\n\s*(?:P\u00e1gina|Page)\s+\d+\s*(?:de|of)\s+\d+\s*\n', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    text = _expand_acronyms(text)
    return text.strip()


def extract_pdf(filepath):
    try:
        import fitz
        text_parts = []
        with fitz.open(str(filepath)) as doc:
            for page in doc:
                t = page.get_text("text")
                if t.strip():
                    text_parts.append(t)
        return clean_text("\n\n".join(text_parts), "pdf")
    except Exception as e:
        logger.error(f"Error PDF {filepath}: {e}")
        return ""


def extract_docx(filepath):
    try:
        from docx import Document
        doc = Document(str(filepath))
        parts = []
        for p in doc.paragraphs:
            if not p.text.strip():
                continue
            if p.style.name.startswith('Heading'):
                parts.append(f"## {p.text.strip()}")
            else:
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        return clean_text("\n\n".join(parts), "docx")
    except Exception as e:
        logger.error(f"Error DOCX {filepath}: {e}")
        return ""


def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = [s.text.strip() for s in slide.shapes
                     if hasattr(s, "text") and s.text.strip()]
            notes = ""
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    notes = notes_frame.text.strip()
            slide_text = f"[Diapositiva {i+1}]\n" + "\n".join(texts)
            if notes:
                slide_text += f"\n[Notas]: {notes}"
            if texts or notes:
                parts.append(slide_text)
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"Error PPTX {filepath}: {e}")
        return ""


def extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(h) if h is not None else "" for h in rows[0]]
            for row in rows[1:]:
                row_text = " | ".join(
                    f"{headers[i]}: {val}"
                    for i, val in enumerate(row)
                    if val is not None and str(val).strip()
                )
                if row_text.strip():
                    parts.append(f"[Hoja: {sheet_name}]\n"
                                 + " | ".join(headers) + "\n" + row_text)
        wb.close()
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"Error XLSX {filepath}: {e}")
        return ""


def extract_csv(filepath):
    try:
        import csv
        with open(str(filepath), encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = list(reader)
        if not rows:
            return ""
        headers = rows[0]
        chunks = []
        for row in rows[1:]:
            row_text = " | ".join(
                f"{headers[i]}: {val}"
                for i, val in enumerate(row)
                if i < len(headers) and val.strip()
            )
            if row_text.strip():
                chunks.append(row_text)
        return "\n\n".join(chunks)
    except Exception as e:
        logger.error(f"Error CSV {filepath}: {e}")
        return ""


def extract_txt(filepath):
    try:
        return Path(filepath).read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logger.error(f"Error TXT {filepath}: {e}")
        return ""


def extract_html(filepath):
    try:
        from bs4 import BeautifulSoup
        html = Path(filepath).read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(html, "lxml")
        for tag in soup(["script", "style", "head", "nav", "footer", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        return clean_text("\n".join(lines), "html")
    except Exception as e:
        logger.error(f"Error HTML {filepath}: {e}")
        return ""
