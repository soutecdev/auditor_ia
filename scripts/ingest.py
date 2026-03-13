import os
import re
import hashlib
import logging
import io
from pathlib import Path
from datetime import datetime

from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(message)s')
logger = logging.getLogger(__name__)

# ── Configuración AuditIA ──────────────────────────────────────────────────────
CORPUS_DIR      = "/data/AI_Projects/AuditIA/corpus"
QDRANT_URL      = "http://localhost:6333"
COLLECTION_NAME = "auditia"                                  # ← separado de "corporativo"
EMBEDDING_MODEL = "nomic-ai/nomic-embed-text-v1.5"
EMBEDDING_CACHE = "/data/AI_Projects/SonIA/models/nomic"     # ← reutiliza el modelo, no duplica
EMBEDDING_DIMS  = 768
CHUNK_SIZE      = 500
CHUNK_OVERLAP   = 50
BATCH_SIZE      = 16

client = QdrantClient(url=QDRANT_URL)
model  = SentenceTransformer(EMBEDDING_MODEL, cache_folder=EMBEDDING_CACHE, trust_remote_code=True)

# ── NORMA_REGISTRY ─────────────────────────────────────────────────────────────
# Detecta metadatos normativos automáticamente por nombre de archivo.
# Añade aquí cualquier norma nueva que vayas a ingestar.
# El campo "keywords" permite detectar variantes del nombre del archivo.
NORMA_REGISTRY = {
    # ── ISA / NIA (Normas Internacionales de Auditoría) ─────────────────────
    "ISA_200": {
        "keywords": ["ISA_200", "ISA200", "NIA_200", "NIA200"],
        "norma_codigo": "ISA-200", "norma_familia": "ISA", "norma_version": "2009",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["general"],
    },
    "ISA_240": {
        "keywords": ["ISA_240", "ISA240", "NIA_240", "NIA240"],
        "norma_codigo": "ISA-240", "norma_familia": "ISA", "norma_version": "2009",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["fraude"],
    },
    "ISA_315": {
        "keywords": ["ISA_315", "ISA315", "NIA_315", "NIA315"],
        "norma_codigo": "ISA-315", "norma_familia": "ISA", "norma_version": "2022",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["riesgos", "control_interno"],
    },
    "ISA_330": {
        "keywords": ["ISA_330", "ISA330", "NIA_330", "NIA330"],
        "norma_codigo": "ISA-330", "norma_familia": "ISA", "norma_version": "2009",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["control_interno"],
    },
    "ISA_500": {
        "keywords": ["ISA_500", "ISA500", "NIA_500", "NIA500"],
        "norma_codigo": "ISA-500", "norma_familia": "ISA", "norma_version": "2009",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["general"],
    },
    "ISA_700": {
        "keywords": ["ISA_700", "ISA700", "NIA_700", "NIA700"],
        "norma_codigo": "ISA-700", "norma_familia": "ISA", "norma_version": "2015",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["general"],
    },

    # ── NIIF / IFRS ──────────────────────────────────────────────────────────
    "NIIF_9": {
        "keywords": ["NIIF_9", "NIIF9", "IFRS_9", "IFRS9"],
        "norma_codigo": "NIIF-9", "norma_familia": "NIIF", "norma_version": "2014",
        "organismo": "IASB", "ambito": "internacional",
        "ciclo_audit": ["tesoreria", "consolidacion"],
    },
    "NIIF_15": {
        "keywords": ["NIIF_15", "NIIF15", "IFRS_15", "IFRS15"],
        "norma_codigo": "NIIF-15", "norma_familia": "NIIF", "norma_version": "2014",
        "organismo": "IASB", "ambito": "internacional",
        "ciclo_audit": ["ventas"],
    },
    "NIIF_16": {
        "keywords": ["NIIF_16", "NIIF16", "IFRS_16", "IFRS16"],
        "norma_codigo": "NIIF-16", "norma_familia": "NIIF", "norma_version": "2016",
        "organismo": "IASB", "ambito": "internacional",
        "ciclo_audit": ["activos_fijos"],
    },

    # ── PCAOB ────────────────────────────────────────────────────────────────
    "PCAOB_AS2201": {
        "keywords": ["PCAOB_AS2201", "AS2201", "AS_2201"],
        "norma_codigo": "PCAOB-AS2201", "norma_familia": "PCAOB", "norma_version": "2007",
        "organismo": "PCAOB", "ambito": "internacional",
        "ciclo_audit": ["control_interno"],
        "industria": ["cotizadas"],
    },
    "PCAOB_AS2301": {
        "keywords": ["PCAOB_AS2301", "AS2301", "AS_2301"],
        "norma_codigo": "PCAOB-AS2301", "norma_familia": "PCAOB", "norma_version": "2010",
        "organismo": "PCAOB", "ambito": "internacional",
        "ciclo_audit": ["control_interno"],
        "industria": ["cotizadas"],
    },

    # ── COSO ─────────────────────────────────────────────────────────────────
    "COSO_2013": {
        "keywords": ["COSO_2013", "COSO2013", "COSO_CI"],
        "norma_codigo": "COSO-2013", "norma_familia": "COSO", "norma_version": "2013",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["control_interno", "riesgos"],
    },
    "COSO_ERM": {
        "keywords": ["COSO_ERM", "COSO_2017", "COSO2017"],
        "norma_codigo": "COSO-ERM-2017", "norma_familia": "COSO", "norma_version": "2017",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["riesgos"],
    },

    # ── SOX ──────────────────────────────────────────────────────────────────
    "SOX": {
        "keywords": ["SOX", "Sarbanes_Oxley", "SarbanesOxley"],
        "norma_codigo": "SOX", "norma_familia": "SOX", "norma_version": "2002",
        "organismo": "US_Congress", "ambito": "internacional",
        "ciclo_audit": ["control_interno", "fraude"],
        "industria": ["cotizadas"],
    },
    "NIIF_MARCO_CONCEPTUAL": {
        "keywords": ["NIIF_MarcoConceptual", "conceptual-framework"],
        "norma_codigo": "NIIF-MARCO", "norma_familia": "NIIF", "norma_version": "2024",
        "organismo": "IASB", "ambito": "internacional",
        "ciclo_audit": ["general"],
    },
    "COSO_2013_EXEC": {
        "keywords": ["COSO_2013_ExecutiveSummary", "COSO_2013"],
        "norma_codigo": "COSO-2013", "norma_familia": "COSO", "norma_version": "2013",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["control_interno", "riesgos"],
    },
    "ISA_HANDBOOK_VOL1": {
        "keywords": ["ISA_Handbook_2022_Vol1"],
        "norma_codigo": "ISA-HANDBOOK-V1", "norma_familia": "ISA", "norma_version": "2022",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["general", "control_interno", "riesgos", "fraude"],
    },
    "ISA_HANDBOOK_VOL2": {
        "keywords": ["ISA_Handbook_2022_Vol2"],
        "norma_codigo": "ISA-HANDBOOK-V2", "norma_familia": "ISA", "norma_version": "2022",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["general"],
    },
    "ISA_HANDBOOK_VOL3": {
        "keywords": ["ISA_Handbook_2022_Vol3"],
        "norma_codigo": "ISA-HANDBOOK-V3", "norma_familia": "ISA", "norma_version": "2022",
        "organismo": "IAASB", "ambito": "internacional",
        "ciclo_audit": ["general"],
    },
    # ── LOTTT (Ley Orgánica del Trabajo Venezuela) ─────────────────────────────
    "LOTTT": {
        "keywords": ["LOTTT", "LOT_VE", "LEY_ORGANICA_TRABAJO", "LEY_TRABAJO_VE",
                     "LOTTT_VE", "LOT2012", "LOTTT2012"],
        "norma_codigo": "LOTTT-2012", "norma_familia": "LOTTT", "norma_version": "2012",
        "organismo": "MPPPST", "ambito": "venezuela",
        "ciclo_audit": ["nomina"], "industria": ["general"],
    },
    # ── COT (Código Orgánico Tributario Venezuela) ──────────────────────────────
    "COT_VE": {
        "keywords": ["COT_VE", "COT_VE", "CODIGO_ORGANICO_TRIBUTARIO", "COT2020"],
        "norma_codigo": "COT-2020", "norma_familia": "COT", "norma_version": "2020",
        "organismo": "SENIAT", "ambito": "venezuela",
        "ciclo_audit": ["impuestos"], "industria": ["general"],
    },
    # ── ISLR (Ley de Impuesto Sobre la Renta Venezuela) ──────────────────────
    "ISLR_VE": {
        "keywords": ["ISLR_VE", "ISLR", "IMPUESTO_SOBRE_LA_RENTA", "ISLR_VE"],
        "norma_codigo": "ISLR-2015", "norma_familia": "ISLR", "norma_version": "2015",
        "organismo": "SENIAT", "ambito": "venezuela",
        "ciclo_audit": ["impuestos"], "industria": ["general"],
    },

    # ── IIA ───────────────────────────────────────────────────────────────────
    "IIA_NORMAS_GLOBALES": {
        "keywords": ["IIA_Normas_Globales", "IIA_NORMAS", "IPPF"],
        "norma_codigo": "IIA-2024", "norma_familia": "IIA", "norma_version": "2024",
        "organismo": "IIA", "ambito": "internacional",
        "ciclo_audit": ["control_interno", "riesgos", "fraude", "tecnologia"],
        "confidencialidad": "publica",
    },
    "IIA_THREE_LINES": {
        "keywords": ["IIA_Three_Lines", "THREE_LINES_MODEL", "Three_Lines"],
        "norma_codigo": "IIA-TLM-2020", "norma_familia": "IIA", "norma_version": "2020",
        "organismo": "IIA", "ambito": "internacional",
        "ciclo_audit": ["control_interno", "riesgos"],
        "confidencialidad": "publica",
    },
    # ── COSO adicionales ──────────────────────────────────────────────────────
    "COSO_FRM": {
        "keywords": ["COSO_FRM", "Fraud_Risk_Management", "COSO_FRAUD", "FRM_Executive"],
        "norma_codigo": "COSO-FRM-2023", "norma_familia": "COSO", "norma_version": "2023",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["fraude", "control_interno", "riesgos"],
        "confidencialidad": "publica",
    },
    "COSO_ICIF": {
        "keywords": ["COSO_ICIF", "Internal_Control_Integrated", "ICIF_2012"],
        "norma_codigo": "COSO-ICIF-2012", "norma_familia": "COSO", "norma_version": "2012",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["control_interno", "riesgos", "fraude"],
        "confidencialidad": "publica",
    },
    "COSO_ERM_2017": {
        "keywords": ["COSO_ERM_2017", "Enterprise_Risk_Management", "ERM_2017"],
        "norma_codigo": "COSO-ERM-2017", "norma_familia": "COSO", "norma_version": "2017",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["riesgos", "control_interno"],
        "confidencialidad": "publica",
    },
    "COSO_ERM_FAQ": {
        "keywords": ["COSO_ERM_FAQ", "ERM_FAQ"],
        "norma_codigo": "COSO-ERM-FAQ", "norma_familia": "COSO", "norma_version": "2017",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["riesgos", "control_interno"],
        "confidencialidad": "publica",
    },
    "COSO_GENAI": {
        "keywords": ["COSO_GenAI", "GenAI_IC", "Generative_AI_Control", "COSO_AI"],
        "norma_codigo": "COSO-GENAI-2026", "norma_familia": "COSO", "norma_version": "2026",
        "organismo": "COSO", "ambito": "internacional",
        "ciclo_audit": ["tecnologia", "control_interno", "riesgos"],
        "confidencialidad": "publica",
    },
    # ── Documentos internos EY ─────────────────────────────────────────────────
    "EY_INTERNO": {
        "keywords": ["EY_Taller", "EY_INTERNO", "EY_"],
        "norma_codigo": None, "norma_familia": "EY_INTERNO", "norma_version": None,
        "organismo": "EY", "ambito": "internacional",
        "ciclo_audit": ["control_interno", "tecnologia", "riesgos"],
        "confidencialidad": "restringida",
    },
}

# Ciclos de auditoría válidos (para referencia)
# compras | ventas | nomina | tesoreria | inventario |
# activos_fijos | impuestos | consolidacion | riesgos |
# control_interno | fraude | tecnologia | sostenibilidad | general


def _detect_norma_metadata(filename: str) -> dict:
    """
    Detecta metadatos normativos automáticamente por nombre de archivo.
    Devuelve dict con campos de auditoría si hay coincidencia,
    o valores por defecto si el archivo no es una norma conocida.
    """
    fname_upper = filename.upper()

    for _key, meta in NORMA_REGISTRY.items():
        for kw in meta["keywords"]:
            if kw.upper() in fname_upper:
                return {
                    "doc_type":         "normativa",
                    "norma_codigo":     meta["norma_codigo"],
                    "norma_familia":    meta["norma_familia"],
                    "norma_version":    meta.get("norma_version", ""),
                    "norma_vigente":    True,
                    "organismo":        meta["organismo"],
                    "ambito":           meta["ambito"],
                    "ciclo_audit":      meta.get("ciclo_audit", ["general"]),
                    "industria":        meta.get("industria", ["general"]),
                    "confidencialidad": meta.get("confidencialidad", "publica"),
                    "cliente_id":       None,
                    "engagement_id":    None,
                    "ejercicio_fiscal": None,
                }

    # Archivo no reconocido como norma → tipo genérico de auditoría
    return {
        "doc_type":         "auditia_doc",
        "norma_codigo":     "",
        "norma_familia":    "",
        "norma_version":    "",
        "norma_vigente":    False,
        "organismo":        "",
        "ambito":           "",
        "ciclo_audit":      ["general"],
        "industria":        ["general"],
        "confidencialidad": "interna",
        "cliente_id":       None,
        "engagement_id":    None,
        "ejercicio_fiscal": None,
    }


def _detect_chunk_flags(text: str) -> dict:
    """
    Detecta si un chunk contiene definiciones, requerimientos o guías de aplicación.
    Útil para filtrar búsquedas por tipo de contenido normativo.
    """
    text_lower = text.lower()

    # Requerimiento: "el auditor debe / shall / deberá"
    es_requerimiento = bool(re.search(
        r'\b(el auditor debe|the auditor shall|deberá|debe evaluar|debe obtener|shall obtain|shall evaluate)\b',
        text_lower
    ))

    # Definición: glosario o sección de definiciones
    es_definicion = bool(re.search(
        r'\b(se define como|significa|definition|definición|a los efectos de|for purposes of)\b',
        text_lower
    ))

    # Guía de aplicación: párrafos A1, A2, A3... o "material de aplicación"
    es_guia_aplicacion = bool(re.search(
        r'\b(a\d+\.|material de aplicación|application material|guidance)\b',
        text_lower
    ))

    return {
        "es_requerimiento":    es_requerimiento,
        "es_definicion":       es_definicion,
        "es_guia_aplicacion":  es_guia_aplicacion,
    }


# ── Setup colección ────────────────────────────────────────────────────────────
def setup_collection():
    existing = [c.name for c in client.get_collections().collections]
    if COLLECTION_NAME not in existing:
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(size=EMBEDDING_DIMS, distance=Distance.COSINE)
        )
        logger.info(f"Coleccion '{COLLECTION_NAME}' creada con {EMBEDDING_DIMS} dims")
    else:
        logger.info(f"Coleccion '{COLLECTION_NAME}' ya existe")


# ── Utilidades ─────────────────────────────────────────────────────────────────
def get_file_hash(filepath):
    h = hashlib.md5()
    with open(filepath, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            h.update(chunk)
    return h.hexdigest()


def _expand_acronyms(text):
    def _add_plain(m):
        original = m.group(0)
        plain = original.replace(".", "")
        if plain == original:
            return original
        return f"{original} ({plain})"
    return re.sub(r'(?:[A-Z]\.){2,}[A-Z]?\.?', _add_plain, text)


def clean_text(text, source_type="generic"):
    if not text or not text.strip():
        return ""
    lines = text.splitlines()
    if len(lines) > 10:
        line_counts = {}
        for line in lines:
            clean_line = line.strip().lower()
            if clean_line and len(clean_line) < 100:
                line_counts[clean_line] = line_counts.get(clean_line, 0) + 1
        page_count = max(1, len(lines) // 40)
        repeated = {l for l, c in line_counts.items() if c > page_count * 0.3 and c > 2}
        lines = [l for l in lines if l.strip().lower() not in repeated]
    text = "\n".join(lines)
    text = re.sub(r'\n\s*\d{1,3}\s*\n', '\n', text)
    text = re.sub(r'\n\s*(?:Página|Page)\s+\d+\s*(?:de|of)\s+\d+\s*\n', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    text = _expand_acronyms(text)
    return text.strip()


def chunk_text(text, max_size=CHUNK_SIZE, overlap_sentences=1):
    if not text or not text.strip():
        return []
    text = text.strip()
    sentences = re.split(r'(?<=[.!?])\s+', text)
    sentences = [s.strip() for s in sentences if s.strip()]
    if len(sentences) <= 1 and len(text) > max_size:
        sentences = [s.strip() for s in text.split('\n') if s.strip()]
    if len(sentences) <= 1 and len(text) > max_size:
        chunks = []
        start = 0
        while start < len(text):
            end = start + max_size
            chunk = text[start:end].strip()
            if len(chunk) >= 30:
                chunks.append(chunk)
            start += max_size - CHUNK_OVERLAP
        return chunks
    chunks = []
    current_chunk = []
    current_len = 0
    for sent in sentences:
        if current_len + len(sent) > max_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = current_chunk[-overlap_sentences:] if overlap_sentences > 0 else []
            current_len = sum(len(s) for s in current_chunk) + len(current_chunk)
        current_chunk.append(sent)
        current_len += len(sent) + 1
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    return [c for c in chunks if len(c) >= 30]


# ── Extractores (sin cambios respecto a SonIA) ─────────────────────────────────
def extract_pdf(filepath):
    try:
        import fitz
        text_parts = []
        page_samples = []
        with fitz.open(str(filepath)) as doc:
            for i, page in enumerate(doc):
                t = page.get_text("text")
                if t.strip():
                    text_parts.append(t)
                if i < 3:
                    page_samples.append(t.strip())
        raw = clean_text("\n\n".join(text_parts), "pdf")
        # Detectar PDF escaneado
        total_sample = " ".join(page_samples)
        if len(total_sample) < 100 or len(raw) < 200:
            logger.info(f"  PDF escaneado detectado: {filepath.name} — usando visión")
            return _extract_pdf_vision_local(filepath)
        return raw
    except Exception as e:
        logger.error(f"Error PDF {filepath.name}: {e}")
        return ""


def _extract_pdf_vision_local(filepath):
    """Convierte páginas de PDF escaneado a imágenes y usa visión — paralelo + resolución optimizada."""
    import fitz, base64, httpx, asyncio

    CONCURRENCY = 3
    MATRIX_SCALE = 1.5

    prompt = """Transcribe todo el texto visible en esta página manteniendo estructura y formato.
Si hay tablas, mantenlas con separadores |. Responde solo con el texto transcrito, sin comentarios."""

    async def _process_page(sem, i, total, img_b64):
        async with sem:
            async with httpx.AsyncClient(timeout=180) as client:
                resp = await client.post("http://localhost:11434/api/chat", json={
                    "model": "sonia-qwen:9b",
                    "think": False,
                    "stream": False,
                    "messages": [{"role": "user", "content": prompt, "images": [img_b64]}]
                })
                text = resp.json().get("message", {}).get("content", "").strip()
                logger.debug(f"  Página {i+1}/{total}: {len(text)} chars")
                return i, text

    async def _process_all(pages_b64):
        sem = asyncio.Semaphore(CONCURRENCY)
        tasks = [_process_page(sem, i, len(pages_b64), b64) for i, b64 in enumerate(pages_b64)]
        results = await asyncio.gather(*tasks)
        return sorted(results, key=lambda x: x[0])

    try:
        pages_b64 = []
        with fitz.open(str(filepath)) as doc:
            total = len(doc)
            logger.info(f"  Procesando {total} páginas con visión paralela (x{CONCURRENCY}): {filepath.name}")
            mat = fitz.Matrix(MATRIX_SCALE, MATRIX_SCALE)
            for page in doc:
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("jpeg")
                pages_b64.append(base64.b64encode(img_bytes).decode("utf-8"))

        results = asyncio.run(_process_all(pages_b64))

        texts = []
        for i, text in results:
            if text:
                texts.append(f"[Página {i+1}/{total}]\n{text}")

        result = "\n\n".join(texts)
        logger.info(f"  Vision PDF OK: {len(result)} chars — {filepath.name}")
        return result
    except Exception as e:
        logger.error(f"  Error vision PDF {filepath.name}: {e}")
        return ""

def extract_pdf_metadata(filepath):
    meta = {"title": filepath.stem, "author": "", "created_date": "", "page_count": 0}
    try:
        import fitz
        doc = fitz.open(str(filepath))
        m = doc.metadata
        if m.get("title"):  meta["title"]        = m["title"]
        if m.get("author"): meta["author"]       = m["author"]
        if m.get("creationDate"): meta["created_date"] = m["creationDate"]
        meta["page_count"] = len(doc)
        doc.close()
    except Exception as e:
        logger.debug(f"  No se pudieron extraer metadatos de {filepath.name}: {e}")
    return meta

def extract_docx(filepath):
    try:
        from docx import Document
        doc = Document(str(filepath))
        parts = []
        for p in doc.paragraphs:
            if not p.text.strip(): continue
            if p.style.name.startswith('Heading'):
                parts.append(f"## {p.text.strip()}")
            else:
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text: parts.append(row_text)
        return clean_text("\n\n".join(parts), "docx")
    except Exception as e:
        logger.error(f"Error DOCX {filepath.name}: {e}")
        return ""

def extract_docx_metadata(filepath):
    meta = {"title": filepath.stem, "author": "", "created_date": ""}
    try:
        from docx import Document
        doc = Document(str(filepath))
        props = doc.core_properties
        if props.title:   meta["title"]        = props.title
        if props.author:  meta["author"]       = props.author
        if props.created: meta["created_date"] = str(props.created)
    except Exception as e:
        logger.debug(f"  No se pudieron extraer metadatos de {filepath.name}: {e}")
    return meta

def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = [s.text.strip() for s in slide.shapes
                     if hasattr(s, "text") and s.text.strip()]
            notes = ""
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    notes = notes_frame.text.strip()
            slide_text = f"[Diapositiva {i+1}]\n" + "\n".join(texts)
            if notes: slide_text += f"\n[Notas]: {notes}"
            if texts or notes: parts.append(slide_text)
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"Error PPTX {filepath.name}: {e}")
        return ""

def extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows: continue
            headers = [str(h) if h is not None else "" for h in rows[0]]
            sheet_chunks = []
            for row in rows[1:]:
                row_text = " | ".join(
                    f"{headers[i]}: {val}"
                    for i, val in enumerate(row)
                    if val is not None and str(val).strip()
                )
                if row_text.strip():
                    sheet_chunks.append(f"[Hoja: {sheet_name}]\n"
                                        + " | ".join(headers) + "\n" + row_text)
            parts.extend(sheet_chunks)
        wb.close()
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"Error XLSX {filepath.name}: {e}")
        return ""

def extract_csv(filepath):
    try:
        import csv
        with open(str(filepath), encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = list(reader)
        if not rows: return ""
        headers = rows[0]
        chunks = []
        for row in rows[1:]:
            row_text = " | ".join(
                f"{headers[i]}: {val}"
                for i, val in enumerate(row)
                if i < len(headers) and val.strip()
            )
            if row_text.strip(): chunks.append(row_text)
        return "\n\n".join(chunks)
    except Exception as e:
        logger.error(f"Error CSV {filepath.name}: {e}")
        return ""

def extract_text(filepath):
    try:
        return filepath.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logger.error(f"Error TXT/MD {filepath.name}: {e}")
        return ""

def extract_html(filepath):
    try:
        from bs4 import BeautifulSoup
        html = filepath.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(html, "lxml")
        for tag in soup(["script", "style", "head", "nav", "footer", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        return clean_text("\n".join(lines), "html")
    except Exception as e:
        logger.error(f"Error HTML {filepath.name}: {e}")
        return ""


# ── Procesado y indexación ─────────────────────────────────────────────────────
def process_file(filepath):
    suffix = filepath.suffix.lower()
    try:
        if suffix == ".pdf":
            text = extract_pdf(filepath)
            chunks = chunk_text(text)
            return [{"text": c, "page": i} for i, c in enumerate(chunks)]
        elif suffix == ".docx":
            text = extract_docx(filepath)
            chunks = chunk_text(text)
            return [{"text": c, "page": i} for i, c in enumerate(chunks)]
        elif suffix == ".pptx":
            text = extract_pptx(filepath)
            chunks = chunk_text(text)
            return [{"text": c, "page": i} for i, c in enumerate(chunks)]
        elif suffix == ".xlsx":
            text = extract_xlsx(filepath)
            lines = [l.strip() for l in text.split("\n\n") if l.strip()]
            return [{"text": l, "page": i} for i, l in enumerate(lines)]
        elif suffix == ".csv":
            text = extract_csv(filepath)
            lines = [l.strip() for l in text.split("\n\n") if l.strip()]
            return [{"text": l, "page": i} for i, l in enumerate(lines)]
        elif suffix in (".html", ".htm"):
            text = extract_html(filepath)
            chunks = chunk_text(text)
            return [{"text": c, "page": i} for i, c in enumerate(chunks)]
        elif suffix in (".txt", ".md"):
            text = extract_text(filepath)
            chunks = chunk_text(text)
            return [{"text": c, "page": i} for i, c in enumerate(chunks)]
        else:
            logger.warning(f"Formato no soportado: {suffix}")
            return []
    except Exception as e:
        logger.error(f"Error procesando {filepath.name}: {e}")
        return []

def get_file_metadata(filepath):
    suffix = filepath.suffix.lower()
    if suffix == ".pdf":   return extract_pdf_metadata(filepath)
    elif suffix == ".docx": return extract_docx_metadata(filepath)
    return {"title": filepath.stem, "author": "", "created_date": ""}

def file_already_indexed(file_hash):
    try:
        from qdrant_client.models import Filter, FieldCondition, MatchValue
        results = client.scroll(
            collection_name=COLLECTION_NAME,
            scroll_filter=Filter(must=[FieldCondition(
                key="checksum", match=MatchValue(value=file_hash)
            )]),
            limit=1
        )
        return len(results[0]) > 0
    except:
        return False

def ingest_file(filepath, force=False):
    file_hash = get_file_hash(str(filepath))
    if not force and file_already_indexed(file_hash):
        logger.info(f"  Sin cambios: {filepath.name}")
        return

    if force:
        try:
            from qdrant_client.models import Filter, FieldCondition, MatchValue
            client.delete(
                collection_name=COLLECTION_NAME,
                points_selector=Filter(must=[FieldCondition(
                    key="source", match=MatchValue(value=filepath.name)
                )])
            )
            logger.info(f"  Chunks anteriores eliminados: {filepath.name}")
        except Exception as e:
            logger.warning(f"  No se pudieron eliminar chunks anteriores: {e}")

    logger.info(f"Procesando: {filepath.name}")
    elements = process_file(filepath)
    if not elements:
        logger.warning(f"  Sin contenido extraido: {filepath.name}")
        return

    file_meta  = get_file_metadata(filepath)
    norma_meta = _detect_norma_metadata(filepath.name)   # ← NUEVO

    # Log si se detectó como norma conocida
    if norma_meta["norma_codigo"]:
        logger.info(f"  Norma detectada: {norma_meta['norma_codigo']} "
                    f"({norma_meta['norma_familia']}) — {norma_meta['ambito']}")
    else:
        logger.info(f"  Documento genérico AuditIA: {filepath.name}")

    logger.info(f"  {len(elements)} chunks extraidos — procesando en lotes para evitar OOM")

    STREAM_BATCH = 200   # chunks por lote: embed + upsert + liberar memoria
    total_indexed = 0

    for start in range(0, len(elements), STREAM_BATCH):
        batch_elems = elements[start:start + STREAM_BATCH]
        texts       = [e["text"] for e in batch_elems]
        embeddings  = model.encode(texts, normalize_embeddings=True,
                                   batch_size=8, show_progress_bar=False)
        points = []
        for j, (emb, elem) in enumerate(zip(embeddings, batch_elems)):
            chunk_flags = _detect_chunk_flags(elem["text"])
            point_id    = abs(hash(f"{file_hash}_{start+j}")) % (2**63)
            points.append(PointStruct(
                id=point_id,
                vector=emb.tolist(),
                payload={
                    "text":               elem["text"],
                    "source":             filepath.name,
                    "file_type":          filepath.suffix.lower().lstrip("."),
                    "page":               elem.get("page", 0),
                    "checksum":           file_hash,
                    "indexed_at":         datetime.now().isoformat(),
                    "title":              file_meta.get("title", filepath.stem),
                    "author":             file_meta.get("author", ""),
                    "doc_type":           norma_meta["doc_type"],
                    "norma_codigo":       norma_meta["norma_codigo"],
                    "norma_familia":      norma_meta["norma_familia"],
                    "norma_version":      norma_meta["norma_version"],
                    "norma_vigente":      norma_meta["norma_vigente"],
                    "organismo":          norma_meta["organismo"],
                    "ambito":             norma_meta["ambito"],
                    "ciclo_audit":        norma_meta["ciclo_audit"],
                    "industria":          norma_meta["industria"],
                    "confidencialidad":   norma_meta["confidencialidad"],
                    "cliente_id":         norma_meta["cliente_id"],
                    "engagement_id":      norma_meta["engagement_id"],
                    "ejercicio_fiscal":   norma_meta["ejercicio_fiscal"],
                    "chunk_index":        start + j,
                    "es_requerimiento":   chunk_flags["es_requerimiento"],
                    "es_definicion":      chunk_flags["es_definicion"],
                    "es_guia_aplicacion": chunk_flags["es_guia_aplicacion"],
                }
            ))
        client.upsert(collection_name=COLLECTION_NAME, points=points)
        total_indexed += len(points)
        logger.info(f"  Lote {start//STREAM_BATCH + 1}: {total_indexed}/{len(elements)} chunks indexados")

    logger.info(f"  Completo: {total_indexed} chunks en colección '{COLLECTION_NAME}'")


def run_ingestion(force=False):
    setup_collection()
    extensions = [".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".txt", ".md", ".html", ".htm"]
    files = [f for ext in extensions
               for f in Path(CORPUS_DIR).rglob(f"*{ext}")]
    if not files:
        logger.warning(f"No se encontraron archivos en {CORPUS_DIR}")
        logger.info(f"Coloca las normativas en: {CORPUS_DIR}")
        return
    logger.info(f"{len(files)} archivos encontrados en corpus AuditIA")
    if force:
        logger.info("Modo FORCE: re-indexando todos los archivos")
    for i, filepath in enumerate(files, 1):
        logger.info(f"[{i}/{len(files)}] {filepath.name}")
        ingest_file(filepath, force=force)
    info = client.get_collection(COLLECTION_NAME)
    logger.info(f"Ingesta completa. Total vectores en '{COLLECTION_NAME}': {info.points_count:,}")


if __name__ == "__main__":
    import sys
    force = "--force" in sys.argv
    run_ingestion(force=force)
